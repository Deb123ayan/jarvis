"""
tools/filesystem.py — File system tools for JARVIS.

Tools: file_read, file_write, file_delete, file_move, file_search, dir_list

Supports reading: .txt / .md / .py / .csv / any plain-text
                  .pdf  (via PyMuPDF  / fitz)
                  .docx (via python-docx)
                  .xlsx (via openpyxl)
                  .pptx (via python-pptx)
"""

import glob
import shutil
from pathlib import Path

from loguru import logger

try:
    import send2trash
    _HAS_SEND2TRASH = True
except ImportError:
    _HAS_SEND2TRASH = False
    logger.warning("send2trash not installed — file_delete will permanently delete files.")

# Optional rich-doc libraries — fail gracefully if not installed
try:
    import fitz  # PyMuPDF
    _HAS_FITZ = True
except ImportError:
    _HAS_FITZ = False
    logger.warning("PyMuPDF (fitz) not installed — PDF reading disabled. Run: pip install pymupdf")

try:
    from docx import Document as _DocxDocument
    _HAS_DOCX = True
except ImportError:
    _HAS_DOCX = False
    logger.warning("python-docx not installed — DOCX reading disabled. Run: pip install python-docx")

try:
    import openpyxl
    _HAS_OPENPYXL = True
except ImportError:
    _HAS_OPENPYXL = False
    logger.warning("openpyxl not installed — XLSX reading disabled. Run: pip install openpyxl")

try:
    from pptx import Presentation as _PptxPresentation
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False
    logger.warning("python-pptx not installed — PPTX reading disabled. Run: pip install python-pptx")

from tools.registry import Tool


# ---------------------------------------------------------------------------
# Handlers
# ---------------------------------------------------------------------------

def _read_pdf(p: Path) -> str:
    """Extract text from a PDF using PyMuPDF."""
    if not _HAS_FITZ:
        return "Error: PyMuPDF is not installed. Run: pip install pymupdf"
    try:
        doc = fitz.open(str(p))
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                pages.append(f"--- Page {i + 1} ---\n{text.strip()}")
        doc.close()
        if not pages:
            return "(PDF appears to be image-only or has no extractable text.)"
        full = "\n\n".join(pages)
        # Cap at ~8000 chars to avoid flooding TTS
        if len(full) > 8000:
            return full[:8000] + f"\n\n[...truncated — {len(full)} total characters]"
        return full
    except Exception as e:
        return f"Error reading PDF: {e}"


def _read_docx(p: Path) -> str:
    """Extract text from a DOCX file using python-docx."""
    if not _HAS_DOCX:
        return "Error: python-docx is not installed. Run: pip install python-docx"
    try:
        doc = _DocxDocument(str(p))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        full = "\n".join(paragraphs)
        if len(full) > 8000:
            return full[:8000] + f"\n\n[...truncated — {len(full)} total characters]"
        return full if full else "(Document appears to be empty.)"
    except Exception as e:
        return f"Error reading DOCX: {e}"


def _read_xlsx(p: Path) -> str:
    """Extract text from an XLSX spreadsheet using openpyxl."""
    if not _HAS_OPENPYXL:
        return "Error: openpyxl is not installed. Run: pip install openpyxl"
    try:
        wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(c) if c is not None else "" for c in row)
                if row_text.strip():
                    lines.append(row_text)
        wb.close()
        full = "\n".join(lines)
        if len(full) > 8000:
            return full[:8000] + f"\n\n[...truncated — {len(full)} total characters]"
        return full if full else "(Spreadsheet appears to be empty.)"
    except Exception as e:
        return f"Error reading XLSX: {e}"


def _read_pptx(p: Path) -> str:
    """Extract text from a PPTX presentation using python-pptx."""
    if not _HAS_PPTX:
        return "Error: python-pptx is not installed. Run: pip install python-pptx"
    try:
        prs = _PptxPresentation(str(p))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
        full = "\n\n".join(slides)
        if len(full) > 8000:
            return full[:8000] + f"\n\n[...truncated — {len(full)} total characters]"
        return full if full else "(Presentation appears to be empty.)"
    except Exception as e:
        return f"Error reading PPTX: {e}"


# Map of extensions → reader functions
_RICH_READERS = {
    ".pdf":  _read_pdf,
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".pptx": _read_pptx,
}


def file_read(path: str) -> str:
    """Read and return the contents of a file. Supports PDF, DOCX, XLSX, PPTX, and plain text."""
    p = Path(path).expanduser()
    if not p.exists():
        return f"Error: File not found: {path}"
    if not p.is_file():
        return f"Error: Not a file: {path}"

    ext = p.suffix.lower()

    # Route to rich-doc readers for known binary formats
    if ext in _RICH_READERS:
        return _RICH_READERS[ext](p)

    # Plain-text fallback
    try:
        size = p.stat().st_size
        if size > 5_000_000:  # 5 MB safety limit for plain text
            return f"Error: File is too large to read ({size // 1024} KB). Use a text editor."
        return p.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        return f"Error reading file: {e}"


def file_write(path: str, content: str) -> str:
    """Write content to a file, creating parent directories if needed."""
    p = Path(path).expanduser()
    try:
        p.parent.mkdir(parents=True, exist_ok=True)
        p.write_text(content, encoding="utf-8")
        return f"File written successfully: {path} ({len(content)} characters)"
    except Exception as e:
        return f"Error writing file: {e}"


def file_delete(path: str) -> str:
    """Move a file to the recycle bin (safe delete via send2trash)."""
    p = Path(path).expanduser()
    if not p.exists():
        return f"Error: File not found: {path}"
    try:
        if _HAS_SEND2TRASH:
            send2trash.send2trash(str(p))
            return f"File moved to recycle bin: {path}"
        else:
            p.unlink()
            return f"File permanently deleted: {path}"
    except Exception as e:
        return f"Error deleting file: {e}"


def file_move(source: str, destination: str) -> str:
    """Move or rename a file/directory."""
    src = Path(source).expanduser()
    dst = Path(destination).expanduser()
    if not src.exists():
        return f"Error: Source not found: {source}"
    try:
        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(src), str(dst))
        return f"Moved: {source} → {destination}"
    except Exception as e:
        return f"Error moving file: {e}"


def file_search(directory: str, pattern: str) -> str:
    """Search for files matching a glob pattern inside a directory."""
    base = Path(directory).expanduser()
    if not base.is_dir():
        return f"Error: Not a directory: {directory}"
    try:
        matches = list(base.glob(f"**/{pattern}"))
        if not matches:
            return f"No files found matching '{pattern}' in {directory}"
        # Limit output to first 50 results
        lines = [str(m) for m in matches[:50]]
        suffix = f"\n... and {len(matches) - 50} more." if len(matches) > 50 else ""
        return "\n".join(lines) + suffix
    except Exception as e:
        return f"Error searching: {e}"


def dir_list(path: str) -> str:
    """List the contents of a directory."""
    p = Path(path).expanduser()
    if not p.exists():
        return f"Error: Path not found: {path}"
    if not p.is_dir():
        return f"Error: Not a directory: {path}"
    try:
        items = sorted(p.iterdir(), key=lambda x: (x.is_file(), x.name.lower()))
        lines = []
        for item in items[:100]:
            kind = "DIR " if item.is_dir() else "FILE"
            size = f"({item.stat().st_size:,} bytes)" if item.is_file() else ""
            lines.append(f"[{kind}] {item.name} {size}")
        if len(list(p.iterdir())) > 100:
            lines.append("... (showing first 100 items)")
        return "\n".join(lines) if lines else "Directory is empty."
    except Exception as e:
        return f"Error listing directory: {e}"


# ---------------------------------------------------------------------------
# Tool definitions
# ---------------------------------------------------------------------------

FILESYSTEM_TOOLS = [
    Tool(
        name="file_read",
        description=(
            "Read the full text contents of a file at the given path. "
            "Supports plain text, PDF (.pdf), Word documents (.docx), "
            "Excel spreadsheets (.xlsx), and PowerPoint presentations (.pptx). "
            "Use when the user wants to view, summarize, or reference a file's contents."
        ),
        args_schema={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Absolute or home-relative path to the file (PDF, DOCX, XLSX, PPTX, or plain text)."}
            },
            "required": ["path"]
        },
        handler=file_read,
        risk_level="low",
    ),
    Tool(
        name="file_write",
        description="Write text content to a file, creating it if it doesn't exist. Overwrites existing content. Use when the user wants to create or save a file.",
        args_schema={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Absolute or home-relative path to the file."},
                "content": {"type": "string", "description": "Text content to write to the file."}
            },
            "required": ["path", "content"]
        },
        handler=file_write,
        risk_level="medium",
    ),
    Tool(
        name="file_delete",
        description="Move a file to the recycle bin. Requires voice confirmation. Use when the user explicitly asks to delete a specific file.",
        args_schema={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Absolute path to the file to delete."}
            },
            "required": ["path"]
        },
        handler=file_delete,
        risk_level="high",
    ),
    Tool(
        name="file_move",
        description="Move or rename a file or directory from source to destination. Requires confirmation. Use when the user wants to move or rename a file.",
        args_schema={
            "type": "object",
            "properties": {
                "source": {"type": "string", "description": "Current absolute path of the file or directory."},
                "destination": {"type": "string", "description": "New absolute path for the file or directory."}
            },
            "required": ["source", "destination"]
        },
        handler=file_move,
        risk_level="high",
    ),
    Tool(
        name="file_search",
        description="Search for files matching a pattern inside a directory. Returns matching file paths. Use when the user asks to find files.",
        args_schema={
            "type": "object",
            "properties": {
                "directory": {"type": "string", "description": "Root directory to search within."},
                "pattern": {"type": "string", "description": "Glob pattern, e.g. '*.py', 'report*.pdf'."}
            },
            "required": ["directory", "pattern"]
        },
        handler=file_search,
        risk_level="low",
    ),
    Tool(
        name="dir_list",
        description="List the files and folders inside a directory. Use when the user asks what files are in a folder.",
        args_schema={
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Absolute path of the directory to list."}
            },
            "required": ["path"]
        },
        handler=dir_list,
        risk_level="low",
    ),
]
